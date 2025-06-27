#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Módulo de leitura de arquivos com suporte a múltiplos formatos
Inclui capacidade de OCR para PDFs escaneados
"""

import os
import re
import json
import xml.etree.ElementTree as ET
import zipfile
import yaml
from pathlib import Path
import pandas as pd
import pdfplumber
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import extract_msg
import eml_parser
import pytesseract
from PIL import Image
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
import io
import email
from email import policy

def extrair_texto(caminho_arquivo):
    """
    Extrai texto de diferentes formatos de arquivo
    
    Args:
        caminho_arquivo (str): Caminho do arquivo a ser processado
        
    Returns:
        str: Texto extraído do arquivo
    """
    try:
        arquivo_path = Path(caminho_arquivo)
        extensao = arquivo_path.suffix.lower()
        
        print(f"  📖 Lendo arquivo: {arquivo_path.name}")
        
        # Documentos de texto
        if extensao == '.txt':
            return extrair_texto_txt(caminho_arquivo)
        elif extensao == '.md':
            return extrair_texto_txt(caminho_arquivo)  # Markdown como texto simples
        elif extensao == '.log':
            return extrair_texto_txt(caminho_arquivo)  # Logs como texto simples
        
        # Documentos PDF
        elif extensao == '.pdf':
            return extrair_texto_pdf(caminho_arquivo)
        
        # Documentos Office
        elif extensao == '.docx':
            return extrair_texto_docx(caminho_arquivo)
        elif extensao == '.xlsx':
            return extrair_texto_xlsx(caminho_arquivo)
        elif extensao == '.pptx':
            return extrair_texto_pptx(caminho_arquivo)
        elif extensao == '.csv':
            return extrair_texto_csv(caminho_arquivo)
        
        # E-mails
        elif extensao == '.msg':
            return extrair_texto_msg(caminho_arquivo)
        elif extensao == '.eml':
            return extrair_texto_eml(caminho_arquivo)
        
        # Imagens com OCR
        elif extensao in ['.jpg', '.jpeg', '.png']:
            return extrair_texto_imagem(caminho_arquivo)
        
        # Arquivos estruturados
        elif extensao == '.xml':
            return extrair_texto_xml(caminho_arquivo)
        elif extensao == '.json':
            return extrair_texto_json(caminho_arquivo)
        elif extensao in ['.yaml', '.yml']:
            return extrair_texto_yaml(caminho_arquivo)
        elif extensao == '.html':
            return extrair_texto_html(caminho_arquivo)
        elif extensao == '.rtf':
            return extrair_texto_rtf(caminho_arquivo)
        
        # Arquivos compactados
        elif extensao == '.zip':
            return extrair_texto_zip(caminho_arquivo)
        
        else:
            print(f"  ⚠️  Formato não suportado: {extensao}")
            return ""
            
    except Exception as e:
        print(f"  ❌ Erro ao extrair texto: {str(e)}")
        return ""

def extrair_texto_txt(caminho_arquivo):
    """
    Extrai texto de arquivo .txt
    """
    try:
        # Tentar diferentes encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(caminho_arquivo, 'r', encoding=encoding) as arquivo:
                    return arquivo.read()
            except UnicodeDecodeError:
                continue
        
        # Se nenhum encoding funcionou, usar bytes
        with open(caminho_arquivo, 'rb') as arquivo:
            conteudo = arquivo.read()
            return conteudo.decode('utf-8', errors='ignore')
            
    except Exception as e:
        print(f"    ❌ Erro ao ler arquivo TXT: {str(e)}")
        return ""

def extrair_texto_pdf(caminho_arquivo):
    """
    Extrai texto de PDF usando pdfplumber e OCR como fallback
    """
    texto = ""
    
    try:
        # Tentar extração de texto normal primeiro
        with pdfplumber.open(caminho_arquivo) as pdf:
            for i, pagina in enumerate(pdf.pages):
                try:
                    texto_pagina = pagina.extract_text()
                    if texto_pagina:
                        texto += f"\n--- Página {i+1} ---\n"
                        texto += texto_pagina
                        texto += "\n"
                except Exception as e:
                    print(f"    ⚠️  Erro na página {i+1}: {str(e)}")
                    continue
        
        # Se não extraiu texto suficiente, tentar OCR
        if len(texto.strip()) < 50:  # Texto muito curto pode indicar PDF escaneado
            print("    🔍 Texto insuficiente, tentando OCR...")
            texto_ocr = extrair_texto_pdf_ocr(caminho_arquivo)
            if texto_ocr:
                texto = texto_ocr
                
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do PDF: {str(e)}")
        
    return texto

def extrair_texto_pdf_ocr(caminho_arquivo):
    """
    Extrai texto de PDF usando OCR (pytesseract)
    """
    try:
        import fitz  # PyMuPDF - fallback se não tiver, usar pdfplumber
        
        texto_ocr = ""
        doc = fitz.open(caminho_arquivo)
        
        for num_pagina in range(len(doc)):
            pagina = doc.load_page(num_pagina)
            pix = pagina.get_pixmap()
            img_data = pix.tobytes("png")
            
            # Converter para PIL Image
            img = Image.open(io.BytesIO(img_data))
            
            # Aplicar OCR
            texto_pagina = pytesseract.image_to_string(img, lang='por')
            
            if texto_pagina.strip():
                texto_ocr += f"\n--- Página {num_pagina+1} (OCR) ---\n"
                texto_ocr += texto_pagina
                texto_ocr += "\n"
        
        doc.close()
        return texto_ocr
        
    except ImportError:
        print("    ⚠️  PyMuPDF não disponível para OCR")
        return ""
    except Exception as e:
        print(f"    ❌ Erro no OCR: {str(e)}")
        return ""

def extrair_texto_docx(caminho_arquivo):
    """
    Extrai texto de documento Word (.docx)
    """
    try:
        doc = Document(caminho_arquivo)
        texto = ""
        
        # Extrair texto dos parágrafos
        for paragrafo in doc.paragraphs:
            if paragrafo.text.strip():
                texto += paragrafo.text + "\n"
        
        # Extrair texto das tabelas
        for tabela in doc.tables:
            for linha in tabela.rows:
                for celula in linha.cells:
                    if celula.text.strip():
                        texto += celula.text + " "
                texto += "\n"
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do DOCX: {str(e)}")
        return ""

def extrair_texto_xlsx(caminho_arquivo):
    """
    Extrai texto de planilha Excel (.xlsx)
    """
    try:
        # Ler todas as abas da planilha
        xls = pd.ExcelFile(caminho_arquivo)
        texto = ""
        
        for nome_aba in xls.sheet_names:
            df = pd.read_excel(caminho_arquivo, sheet_name=nome_aba)
            
            texto += f"\n--- Aba: {nome_aba} ---\n"
            
            # Converter DataFrame para texto
            texto += df.to_string(index=False, na_rep='')
            texto += "\n"
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do XLSX: {str(e)}")
        return ""

def extrair_texto_csv(caminho_arquivo):
    """
    Extrai texto de arquivo CSV
    """
    try:
        # Tentar diferentes separadores e encodings
        separadores = [',', ';', '\t']
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            for sep in separadores:
                try:
                    df = pd.read_csv(caminho_arquivo, sep=sep, encoding=encoding)
                    
                    # Se conseguiu ler e tem mais de uma coluna, provavelmente acertou
                    if len(df.columns) > 1:
                        return df.to_string(index=False, na_rep='')
                        
                except Exception:
                    continue
        
        # Se nenhum separador funcionou, ler como texto simples
        return extrair_texto_txt(caminho_arquivo)
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do CSV: {str(e)}")
        return ""

def extrair_texto_msg(caminho_arquivo):
    """
    Extrai texto de e-mail do Outlook (.msg)
    """
    try:
        msg = extract_msg.openMsg(caminho_arquivo)
        
        texto = ""
        
        # Informações básicas do e-mail
        if msg.sender:
            texto += f"De: {msg.sender}\n"
        if msg.to:
            texto += f"Para: {msg.to}\n"
        if msg.subject:
            texto += f"Assunto: {msg.subject}\n"
        if msg.date:
            texto += f"Data: {msg.date}\n"
        
        texto += "\n--- Corpo do E-mail ---\n"
        
        # Corpo do e-mail
        if msg.body:
            texto += msg.body
        
        # Anexos (apenas nomes)
        if hasattr(msg, 'attachments') and msg.attachments:
            texto += "\n--- Anexos ---\n"
            for anexo in msg.attachments:
                if hasattr(anexo, 'longFilename') and anexo.longFilename:
                    texto += f"Anexo: {anexo.longFilename}\n"
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do MSG: {str(e)}")
        return ""

def extrair_texto_pptx(caminho_arquivo):
    """
    Extrai texto de apresentações PowerPoint (.pptx)
    """
    try:
        prs = Presentation(caminho_arquivo)
        texto = ""
        
        for i, slide in enumerate(prs.slides, 1):
            texto += f"\n--- Slide {i} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texto += shape.text + "\n"
                    
                # Verificar tabelas
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text.strip())
                        if row_text:
                            texto += " | ".join(row_text) + "\n"
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do PPTX: {str(e)}")
        return ""

def extrair_texto_eml(caminho_arquivo):
    """
    Extrai texto de e-mails padrão (.eml)
    """
    try:
        with open(caminho_arquivo, 'rb') as f:
            raw_email = f.read()
        
        # Usar eml_parser para análise completa
        ep = eml_parser.EmlParser()
        parsed_email = ep.decode_email_bytes(raw_email)
        
        texto = ""
        
        # Cabeçalhos básicos
        if 'header' in parsed_email:
            header = parsed_email['header']
            if 'from' in header:
                texto += f"De: {header['from']}\n"
            if 'to' in header:
                texto += f"Para: {header['to']}\n"
            if 'subject' in header:
                texto += f"Assunto: {header['subject']}\n"
            if 'date' in header:
                texto += f"Data: {header['date']}\n"
        
        texto += "\n--- Corpo do E-mail ---\n"
        
        # Corpo do e-mail
        if 'body' in parsed_email:
            for body_part in parsed_email['body']:
                if 'content' in body_part:
                    texto += body_part['content'] + "\n"
        
        # Anexos
        if 'attachment' in parsed_email and parsed_email['attachment']:
            texto += "\n--- Anexos ---\n"
            for anexo in parsed_email['attachment']:
                if 'filename' in anexo:
                    texto += f"Anexo: {anexo['filename']}\n"
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do EML: {str(e)}")
        # Fallback para parser básico do Python
        try:
            with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
                msg = email.message_from_file(f, policy=policy.default)
            
            texto = ""
            texto += f"De: {msg.get('From', '')}\n"
            texto += f"Para: {msg.get('To', '')}\n"
            texto += f"Assunto: {msg.get('Subject', '')}\n"
            texto += f"Data: {msg.get('Date', '')}\n\n"
            
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        texto += part.get_content()
            else:
                texto += msg.get_content()
            
            return texto
        except Exception:
            return ""

def extrair_texto_imagem(caminho_arquivo):
    """
    Extrai texto de imagens usando OCR (.jpg, .jpeg, .png)
    """
    try:
        image = Image.open(caminho_arquivo)
        texto = pytesseract.image_to_string(image, lang='por')
        
        if not texto.strip():
            # Tentar com inglês se português falhar
            texto = pytesseract.image_to_string(image, lang='eng')
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto da imagem: {str(e)}")
        return ""

def extrair_texto_xml(caminho_arquivo):
    """
    Extrai texto de arquivos XML
    """
    try:
        tree = ET.parse(caminho_arquivo)
        root = tree.getroot()
        
        def extrair_texto_elemento(elemento):
            texto = ""
            if elemento.text:
                texto += elemento.text.strip() + " "
            for child in elemento:
                texto += extrair_texto_elemento(child)
            if elemento.tail:
                texto += elemento.tail.strip() + " "
            return texto
        
        return extrair_texto_elemento(root)
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do XML: {str(e)}")
        # Fallback para BeautifulSoup
        try:
            with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            soup = BeautifulSoup(content, 'xml')
            return soup.get_text(separator=' ', strip=True)
        except Exception:
            return ""

def extrair_texto_json(caminho_arquivo):
    """
    Extrai texto de arquivos JSON
    """
    try:
        with open(caminho_arquivo, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        def extrair_valores(obj, nivel=0):
            texto = ""
            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, str):
                        texto += f"{key}: {value}\n"
                    else:
                        texto += f"{key}:\n{extrair_valores(value, nivel+1)}"
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    texto += f"[{i}]: {extrair_valores(item, nivel+1)}"
            elif isinstance(obj, str):
                texto += obj + "\n"
            elif obj is not None:
                texto += str(obj) + "\n"
            return texto
        
        return extrair_valores(data)
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do JSON: {str(e)}")
        return ""

def extrair_texto_yaml(caminho_arquivo):
    """
    Extrai texto de arquivos YAML (.yaml, .yml)
    """
    try:
        with open(caminho_arquivo, 'r', encoding='utf-8') as f:
            data = yaml.safe_load(f)
        
        def extrair_valores_yaml(obj, nivel=0):
            texto = ""
            indent = "  " * nivel
            if isinstance(obj, dict):
                for key, value in obj.items():
                    if isinstance(value, str):
                        texto += f"{indent}{key}: {value}\n"
                    else:
                        texto += f"{indent}{key}:\n{extrair_valores_yaml(value, nivel+1)}"
            elif isinstance(obj, list):
                for item in obj:
                    texto += f"{indent}- {extrair_valores_yaml(item, nivel+1)}"
            elif isinstance(obj, str):
                texto += obj + "\n"
            elif obj is not None:
                texto += str(obj) + "\n"
            return texto
        
        return extrair_valores_yaml(data)
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do YAML: {str(e)}")
        return ""

def extrair_texto_html(caminho_arquivo):
    """
    Extrai texto de arquivos HTML
    """
    try:
        with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remover scripts e estilos
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Extrair texto
        texto = soup.get_text(separator='\n', strip=True)
        
        # Limpar quebras de linha excessivas
        texto = re.sub(r'\n\s*\n', '\n\n', texto)
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do HTML: {str(e)}")
        return ""

def extrair_texto_rtf(caminho_arquivo):
    """
    Extrai texto de arquivos RTF
    """
    try:
        with open(caminho_arquivo, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
        
        texto = rtf_to_text(rtf_content)
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do RTF: {str(e)}")
        return ""

def extrair_texto_zip(caminho_arquivo):
    """
    Extrai texto de arquivos dentro de ZIPs
    """
    try:
        texto = ""
        
        with zipfile.ZipFile(caminho_arquivo, 'r') as zip_ref:
            # Listar arquivos no ZIP
            texto += "--- Conteúdo do arquivo ZIP ---\n"
            for file_info in zip_ref.filelist:
                texto += f"Arquivo: {file_info.filename}\n"
                
                # Tentar extrair texto de alguns arquivos
                if file_info.filename.lower().endswith(('.txt', '.csv', '.json', '.xml', '.html')):
                    try:
                        with zip_ref.open(file_info.filename) as f:
                            content = f.read().decode('utf-8', errors='ignore')
                            texto += f"\n--- Conteúdo de {file_info.filename} ---\n"
                            texto += content[:1000]  # Limitar a 1000 caracteres por arquivo
                            if len(content) > 1000:
                                texto += "\n... (conteúdo truncado)\n"
                            texto += "\n"
                    except Exception:
                        continue
        
        return texto
        
    except Exception as e:
        print(f"    ❌ Erro ao extrair texto do ZIP: {str(e)}")
        return ""

def limpar_texto(texto):
    """
    Limpa e normaliza o texto extraído
    """
    if not texto:
        return ""
    
    # Remover caracteres de controle
    texto = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', texto)
    
    # Normalizar quebras de linha
    texto = re.sub(r'\r\n', '\n', texto)
    texto = re.sub(r'\r', '\n', texto)
    
    # Remover linhas vazias excessivas
    texto = re.sub(r'\n\s*\n\s*\n', '\n\n', texto)
    
    # Normalizar espaços
    texto = re.sub(r'[ \t]+', ' ', texto)
    
    return texto.strip()

# Função de teste para o módulo
if __name__ == "__main__":
    print("=== TESTE DO LEITOR DE ARQUIVOS ===")
    
    # Testar com arquivo de exemplo se existir
    teste_arquivo = "data/exemplo_contrato.txt"
    
    if os.path.exists(teste_arquivo):
        print(f"Testando com: {teste_arquivo}")
        texto = extrair_texto(teste_arquivo)
        texto_limpo = limpar_texto(texto)
        
        print(f"Texto extraído ({len(texto_limpo)} caracteres):")
        print("=" * 50)
        print(texto_limpo[:500] + "..." if len(texto_limpo) > 500 else texto_limpo)
        print("=" * 50)
    else:
        print(f"Arquivo de teste não encontrado: {teste_arquivo}")
        print("Crie arquivos na pasta 'data' para testar a extração.")
